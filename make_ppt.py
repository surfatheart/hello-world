# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- palette ----------
DEEP   = RGBColor(0x0B, 0x3D, 0x3B)   # deep teal
ACCENT = RGBColor(0x1E, 0x88, 0x6A)   # green
GOLD   = RGBColor(0xC8, 0x9B, 0x3C)   # gold
INK    = RGBColor(0x22, 0x2A, 0x2E)
GREY   = RGBColor(0x4A, 0x55, 0x59)
LIGHT  = RGBColor(0xF2, 0xF6, 0xF4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

CN = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

def rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
    s.shadow.inherit = False
    return s

def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.12):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, fnt) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = fnt or CN
    return tb

def R(t, sz, col=INK, bold=False, fnt=None):
    return (t, sz, col, bold, fnt)

# ============================================================
# Slide 1 — Title / overview
# ============================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SW, SH, DEEP)
rect(s, 0, 0, Inches(0.28), SH, GOLD)
rect(s, 0, Inches(4.4), SW, Inches(0.04), ACCENT)

txt(s, Inches(0.9), Inches(1.0), Inches(11.6), Inches(0.5),
    [[R("中医药 · 人工智能传承专题", 18, GOLD, True)]])

txt(s, Inches(0.85), Inches(1.65), Inches(11.8), Inches(1.9),
    [[R("名老中医临床辨证思维的", 38, WHITE, True)],
     [R("算法建模与数字孪生", 38, WHITE, True)]],
    line_spacing=1.05)

txt(s, Inches(0.9), Inches(3.55), Inches(11.6), Inches(0.8),
    [[R("核心命题：", 18, GOLD, True),
      R("把口授心传的辨证思维  外显化、可计算化、可大规模部署。", 18, LIGHT, False)]])

txt(s, Inches(0.9), Inches(4.7), Inches(11.6), Inches(2.6),
    [[R("战略价值 · 三个维度", 17, GOLD, True)],
     [R("① 个体经验 → 公共知识资产，打破名医资源高度集中的格局；", 16.5, LIGHT, False)],
     [R("② 以算法建模验证、提炼经验规律，推动中医理论的现代化表达；", 16.5, LIGHT, False)],
     [R("③ 数字化“活态传承”，让名医经验在基层、远程与健康管理中持续发挥作用。", 16.5, LIGHT, False)],
     [R("全球态势：", 16.5, GOLD, True),
      R("中、韩、印等传统医学大国已形成各具特色路径，领域整体由概念探索进入系统落地。", 16.5, LIGHT, False)]],
    space_after=9, line_spacing=1.12)

# ============================================================
# content slide helper
# ============================================================
def content_slide(no, kicker, title, bullets, takeaway):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, SW, Inches(1.5), DEEP)
    rect(s, 0, 0, Inches(0.28), SH, GOLD)
    rect(s, Inches(0.62), Inches(0.4), Inches(0.7), Inches(0.7), GOLD)
    txt(s, Inches(0.62), Inches(0.4), Inches(0.7), Inches(0.7),
        [[R(no, 30, DEEP, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.5), Inches(0.32), Inches(11), Inches(0.4),
        [[R(kicker, 14, GOLD, True)]])
    txt(s, Inches(1.5), Inches(0.62), Inches(11.4), Inches(0.8),
        [[R(title, 25, WHITE, True)]])
    y = Inches(1.9)
    for head, body in bullets:
        rect(s, Inches(0.7), y + Emu(int(Inches(0.06))), Inches(0.15), Inches(0.15), ACCENT)
        txt(s, Inches(1.0), y, Inches(11.7), Inches(1.2),
            [[R(head + "  ", 16.5, DEEP, True), R(body, 15.5, GREY, False)]],
            line_spacing=1.1, space_after=0)
        lines = max(1, (len(head) + len(body)) // 30 + 1)
        y = y + Inches(0.36) + Inches(0.30) * lines
    bar_y = Inches(6.45)
    rect(s, 0, bar_y, SW, Inches(1.05), LIGHT)
    rect(s, 0, bar_y, Inches(0.28), Inches(1.05), GOLD)
    txt(s, Inches(0.7), bar_y, Inches(12.2), Inches(1.05),
        [[R("亮点：", 15.5, GOLD, True), R(takeaway, 15.5, INK, True)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
    return s

# ============================================================
# Slide 2 — 多模态数据采集
# ============================================================
content_slide(
    "1", "数据基础设施", "多模态数据采集：从验案整理到知识仓库",
    [("规模最大语料库",
      "覆盖 200 余位全国名老中医、近 10 万份验案、400 余万条医学问答对，构成迄今最大名医知识语料库。"),
     ("专病库纵深布局",
      "肺癌、冠心病、肺纤维化、慢性肾病、甲状腺等，单病种数百至逾千例，覆盖已延伸至西部地区。"),
     ("非结构化智能解析",
      "深度学习+正则混合方法，疾病/症状/病机/中药实体抽取精确率、召回率、F1 达 88%–90%；Bert-BiLSTM-CRF 提升舌象、脉象识别。"),
     ("国际四诊客观化",
      "韩国韩医研究院 4 年、23 家机构、近 3000 例四象体质多模态库；印度 Nadi Tarangini 脉诊逾 13 万例、1250 余家诊所部署。")],
    "数据规模需配质量管控——专病纵深 × 综合广度互补，避免“数据沼泽”，筑牢建模地基。"
)

# ============================================================
# Slide 3 — 辨证思维建模
# ============================================================
content_slide(
    "2", "核心算法", "辨证思维建模：从规则系统到深度学习",
    [("跨病种性能突破",
      "肾病辨证 DNN 精准度 97%、F1 95%；肺纤维化 81.22%；围绝经期/不孕症/癌痛集中在 85%–97%。"),
     ("多算法优势互补",
      "SVM 适合中等样本高维分类；BP/CART 兼顾可解释与准确率；贝叶斯网络擅长概率推理与不确定性。"),
     ("图神经网络进阶",
      "把中药-症状-证候-靶点建为知识图谱，图卷积聚合高阶路径+注意力融合，在语义层模拟名医配伍逻辑。"),
     ("现存瓶颈",
      "数据标注质量、证候标准不统一、小样本病种过拟合，仍是规模化应用的主要障碍。")],
    "深度学习已具跨病种复制能力，正从“研究可行”迈向“临床可用”。"
)

# ============================================================
# Slide 4 — 数字孪生与活态传承
# ============================================================
content_slide(
    "3", "系统落地", "数字孪生与活态传承：从模型到系统",
    [("可溯的知识表示",
      "证候分解为“证候要素+药性组合”，症状→治法→处方步步可溯，奠定辨证推理的透明化与可审计性。"),
     ("全场景系统集成",
      "经典/图像/影音/心理行为纳入统一框架，数字孪生仿真模拟多次复诊动态；DeepSeek 等国产大模型作“理解与生成”内核。"),
     ("国际范式参照",
      "印度 Ayush Grid “Yoga Posture AI” 2026 年 2 月公开展示，以视觉驱动实时姿态评估，纳入国家数字门户。"),
     ("场景边界延伸",
      "从诊疗决策走向健康维护、运动指导、生活方式管理，并通过数字公共基础设施实现国家级规模化部署。")],
    "数字孪生不是静态存档，而是随临床输入动态响应、持续迭代的“数字名医”。"
)

out = "名老中医辨证思维算法建模与数字孪生_汇报.pptx"
prs.save(out)
print("SAVED:", out)
